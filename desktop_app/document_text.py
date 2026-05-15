from __future__ import annotations

import html
import os
import re
import shutil
import subprocess
import sys
import tarfile
import tempfile
import zipfile
from pathlib import Path
from typing import Iterable

MAX_DOCUMENT_FILES = 80
MAX_TEXT_PER_FILE = 12_000
MAX_DOCUMENT_TEXT = 60_000

ARCHIVE_SUFFIXES = {
    ".zip",
    ".rar",
    ".7z",
    ".tar",
    ".gz",
    ".tgz",
    ".bz2",
    ".xz",
}
TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".tsv",
    ".xml",
    ".json",
    ".yaml",
    ".yml",
    ".ini",
    ".cfg",
    ".html",
    ".htm",
    ".log",
}


def _is_archive(path: Path) -> bool:
    name = path.name.lower()
    return (
        path.suffix.lower() in ARCHIVE_SUFFIXES
        or name.endswith(".tar.gz")
        or name.endswith(".tar.bz2")
        or name.endswith(".tar.xz")
    )


def _is_relative_to(path: Path, root: Path) -> bool:
    try:
        path.resolve().relative_to(root.resolve())
        return True
    except ValueError:
        return False


def _safe_extract_path(root: Path, member_name: str) -> Path:
    target = (root / member_name).resolve()
    if not _is_relative_to(target, root):
        raise RuntimeError(f"Небезопасный путь внутри архива: {member_name}")
    return target


def _runtime_roots() -> list[Path]:
    roots: list[Path] = []
    if getattr(sys, "frozen", False):
        roots.append(Path(sys.executable).resolve().parent)
    bundle_dir = getattr(sys, "_MEIPASS", None)
    if bundle_dir:
        roots.append(Path(bundle_dir).resolve())
    roots.append(Path(__file__).resolve().parent)
    roots.append(Path.cwd())
    return roots


def _rar_tool_candidates() -> list[tuple[str, Path]]:
    candidates: list[tuple[str, Path]] = []
    for root in _runtime_roots():
        candidates.extend(
            (
                ("SEVENZIP_TOOL", root / "7z.exe"),
                ("SEVENZIP_TOOL", root / "7za.exe"),
                ("UNRAR_TOOL", root / "unrar.exe"),
                ("SEVENZIP_TOOL", root / "tools" / "7zip" / "7z.exe"),
                ("SEVENZIP_TOOL", root / "tools" / "7zip" / "7za.exe"),
                ("UNRAR_TOOL", root / "tools" / "unrar" / "unrar.exe"),
            )
        )
    for env_name in ("ProgramFiles", "ProgramFiles(x86)", "LOCALAPPDATA"):
        base = os.environ.get(env_name)
        if not base:
            continue
        candidates.extend(
            (
                ("SEVENZIP_TOOL", Path(base) / "7-Zip" / "7z.exe"),
                ("SEVENZIP_TOOL", Path(base) / "7-Zip" / "7za.exe"),
                ("UNRAR_TOOL", Path(base) / "WinRAR" / "unrar.exe"),
            )
        )
    return candidates


def _configure_rarfile_tools(rarfile) -> None:
    for attr, candidate in _rar_tool_candidates():
        if candidate.is_file():
            setattr(rarfile, attr, str(candidate))
            if attr == "SEVENZIP_TOOL" and hasattr(rarfile, "SEVENZIP2_TOOL"):
                rarfile.SEVENZIP2_TOOL = str(candidate)
            break
    try:
        rarfile.tool_setup(force=True)
    except rarfile.RarCannotExec as exc:
        raise RuntimeError(
            "Не найден распаковщик RAR. Установите 7-Zip или WinRAR на этот компьютер "
            "либо положите 7z.exe рядом с программой в папку tools\\7zip."
        ) from exc


def _sevenzip_tool() -> Path | None:
    for attr, candidate in _rar_tool_candidates():
        if attr == "SEVENZIP_TOOL" and candidate.is_file():
            return candidate
    return None


def _looks_like_html_or_json_error(path: Path) -> str | None:
    try:
        raw = path.read_bytes()[:4096]
    except Exception:
        return None
    prefix = raw.lstrip()[:200].lower()
    if prefix.startswith((b"<!doctype html", b"<html", b"{")):
        preview = raw.decode("utf-8", errors="replace")
        preview = re.sub(r"\s+", " ", preview).strip()[:500]
        return preview or "похоже на HTML/JSON-ответ вместо архива"
    return None


def _extract_with_7zip(path: Path, target_dir: Path) -> None:
    tool = _sevenzip_tool()
    if tool is None:
        raise RuntimeError("7-Zip не найден для fallback-распаковки архива.")
    target_dir.mkdir(parents=True, exist_ok=True)
    result = subprocess.run(
        [str(tool), "x", "-y", f"-o{target_dir}", str(path)],
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=180,
    )
    if result.returncode != 0:
        detail = (result.stderr or result.stdout or "").strip()[:1200]
        raise RuntimeError(f"7-Zip не смог распаковать архив: {detail}")


def _extract_archive(path: Path, target_dir: Path) -> None:
    suffix = path.suffix.lower()
    name = path.name.lower()
    target_dir.mkdir(parents=True, exist_ok=True)
    fake_response = _looks_like_html_or_json_error(path)
    if fake_response:
        raise RuntimeError(
            "Файл с расширением архива не является архивом. "
            f"ЭТП вернула HTML/JSON-ответ вместо файла: {fake_response}"
        )
    if suffix == ".zip":
        try:
            with zipfile.ZipFile(path) as zf:
                for member in zf.infolist():
                    if member.is_dir():
                        continue
                    out = _safe_extract_path(target_dir, member.filename)
                    out.parent.mkdir(parents=True, exist_ok=True)
                    with zf.open(member) as src, out.open("wb") as dst:
                        shutil.copyfileobj(src, dst)
        except (zipfile.BadZipFile, NotImplementedError, RuntimeError) as exc:
            try:
                _extract_with_7zip(path, target_dir)
            except Exception as sevenzip_exc:
                raise RuntimeError(
                    f"ZIP не распакован стандартным модулем ({exc}) "
                    f"и fallback через 7-Zip тоже не сработал: {sevenzip_exc}"
                ) from exc
        return
    if suffix == ".7z":
        import py7zr

        try:
            with py7zr.SevenZipFile(path, mode="r") as archive:
                archive.extractall(path=target_dir)
        except Exception as exc:
            try:
                _extract_with_7zip(path, target_dir)
            except Exception as sevenzip_exc:
                raise RuntimeError(
                    f"7z не распакован py7zr ({exc}) "
                    f"и fallback через 7-Zip тоже не сработал: {sevenzip_exc}"
                ) from exc
        return
    if suffix == ".rar":
        import rarfile

        try:
            _configure_rarfile_tools(rarfile)
            with rarfile.RarFile(path) as archive:
                archive.extractall(path=target_dir)
        except Exception as exc:
            try:
                _extract_with_7zip(path, target_dir)
            except Exception as sevenzip_exc:
                raise RuntimeError(
                    f"RAR не распакован через rarfile ({exc}) "
                    f"и fallback через 7-Zip тоже не сработал: {sevenzip_exc}"
                ) from exc
        return
    if suffix in {".tar", ".gz", ".tgz", ".bz2", ".xz"} or name.endswith((".tar.gz", ".tar.bz2", ".tar.xz")):
        mode = "r:*"
        with tarfile.open(path, mode) as tf:
            for member in tf.getmembers():
                if not member.isfile():
                    continue
                out = _safe_extract_path(target_dir, member.name)
                out.parent.mkdir(parents=True, exist_ok=True)
                src = tf.extractfile(member)
                if src is not None:
                    with src, out.open("wb") as dst:
                        shutil.copyfileobj(src, dst)
        return
    shutil.unpack_archive(str(path), str(target_dir))


def _read_text_file(path: Path) -> str:
    raw = path.read_bytes()
    for enc in ("utf-8-sig", "utf-8", "cp1251", "latin-1"):
        try:
            return raw.decode(enc, errors="replace")
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace")


def _read_docx(path: Path) -> str:
    from docx import Document

    doc = Document(path)
    parts: list[str] = []
    parts.extend(p.text for p in doc.paragraphs if p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for ws in wb.worksheets[:10]:
            parts.append(f"Лист: {ws.title}")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if values:
                    parts.append(" | ".join(values))
                row_count += 1
                if row_count >= 300:
                    parts.append("[лист обрезан]")
                    break
    finally:
        wb.close()
    return "\n".join(parts)


def _read_xls_via_excel(path: Path) -> str:
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    excel = None
    try:
        excel = win32com.client.DispatchEx("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        wb = excel.Workbooks.Open(str(path), ReadOnly=True)
        try:
            parts: list[str] = []
            for sheet in list(wb.Worksheets)[:10]:
                parts.append(f"Лист: {sheet.Name}")
                used = sheet.UsedRange
                values = used.Value
                if values is None:
                    continue
                if not isinstance(values, tuple):
                    values = ((values,),)
                elif values and not isinstance(values[0], tuple):
                    values = (values,)
                for row_index, row in enumerate(values, start=1):
                    row_values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                    if row_values:
                        parts.append(" | ".join(row_values))
                    if row_index >= 300:
                        parts.append("[лист обрезан]")
                        break
            return "\n".join(parts)
        finally:
            wb.Close(False)
    finally:
        if excel is not None:
            excel.Quit()
        pythoncom.CoUninitialize()


def _read_pdf(path: Path) -> str:
    parts: list[str] = []
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        for page in list(reader.pages)[:30]:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                parts.append("")
    except Exception:
        pass
    return "\n".join(p for p in parts if p.strip())


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"Слайд {idx}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text or "").strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


def _read_odf(path: Path) -> str:
    from odf import teletype
    from odf.opendocument import load
    from odf.text import P
    from odf.table import TableCell

    doc = load(str(path))
    parts: list[str] = []
    for node in doc.getElementsByType(P):
        text = teletype.extractText(node).strip()
        if text:
            parts.append(text)
    for cell in doc.getElementsByType(TableCell):
        text = teletype.extractText(cell).strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _read_rtf(path: Path) -> str:
    text = _read_text_file(path)
    text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+\d* ?", " ", text)
    text = text.replace("{", " ").replace("}", " ")
    return re.sub(r"\s+", " ", text).strip()


def _read_doc_via_word(path: Path) -> str:
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    word = None
    try:
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        doc = word.Documents.Open(str(path), ReadOnly=True, AddToRecentFiles=False)
        try:
            return str(doc.Content.Text or "")
        finally:
            doc.Close(False)
    finally:
        if word is not None:
            word.Quit()
        pythoncom.CoUninitialize()


def _extract_text_from_file(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return _read_docx(path)
    if suffix in {".xlsx", ".xlsm"}:
        return _read_xlsx(path)
    if suffix == ".xls":
        return _read_xls_via_excel(path)
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix in {".odt", ".ods", ".odp"}:
        return _read_odf(path)
    if suffix == ".rtf":
        return _read_rtf(path)
    if suffix == ".doc":
        return _read_doc_via_word(path)
    if suffix in TEXT_SUFFIXES:
        return html.unescape(_read_text_file(path))
    return ""


def _walk_files(paths: Iterable[Path]) -> list[Path]:
    files: list[Path] = []
    for path in paths:
        if path.is_file():
            files.append(path)
        elif path.is_dir():
            files.extend(p for p in path.rglob("*") if p.is_file())
    return files


def _unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem, suffix = path.stem, path.suffix
    n = 2
    while True:
        candidate = path.with_name(f"{stem}_{n}{suffix}")
        if not candidate.exists():
            return candidate
        n += 1


def _read_documents_text_from_files(
    files: list[Path],
    progress=None,
    issues: list[dict] | None = None,
    registry: str = "",
) -> str:
    sections: list[str] = []
    for index, path in enumerate(files[:MAX_DOCUMENT_FILES], start=1):
        if progress:
            progress(f"Извлекаю текст документа {index}: {path.name}")
        if _is_archive(path):
            sections.append(f"--- Архив: {path.name} ---\n[архив сохранён; содержимое распаковано рядом]")
            continue
        try:
            text = _extract_text_from_file(path).strip()
        except Exception as e:
            sections.append(f"--- Файл: {path.name} ---\n[не удалось извлечь текст: {e}]")
            if issues is not None:
                issues.append(
                    {
                        "severity": "important",
                        "registry": registry,
                        "file": path.name,
                        "message": f"Не удалось извлечь текст из файла: {e}",
                    }
                )
            continue
        if not text:
            sections.append(f"--- Файл: {path.name} ---\n[текст не извлечён или файл не поддерживается]")
            if issues is not None:
                issues.append(
                    {
                        "severity": "important",
                        "registry": registry,
                        "file": path.name,
                        "message": "Текст не извлечён или формат файла не поддерживается.",
                    }
                )
            continue
        if len(text) > MAX_TEXT_PER_FILE:
            text = text[:MAX_TEXT_PER_FILE] + "\n[текст файла обрезан]"
        sections.append(f"--- Файл: {path.name} ---\n{text}")

    result = "\n\n".join(sections)
    if len(result) > MAX_DOCUMENT_TEXT:
        result = result[:MAX_DOCUMENT_TEXT] + "\n\n[общий текст документов обрезан]"
    return result


def prepare_documents_for_analysis(
    files: list[Path],
    output_dir: Path,
    progress=None,
    issues: list[dict] | None = None,
    registry: str = "",
) -> tuple[str, Path]:
    """Сохраняет разархивированные документы в output_dir и читает все файлы рекурсивно."""
    output_dir.parent.mkdir(parents=True, exist_ok=True)
    staging_dir = Path(
        tempfile.mkdtemp(prefix=f"{output_dir.name}_source_", dir=output_dir.parent)
    )
    staged_files: list[Path] = []
    for source in files:
        source = Path(source)
        try:
            if not source.is_file():
                if issues is not None:
                    issues.append(
                        {
                            "severity": "critical",
                            "registry": registry,
                            "file": source.name,
                            "message": f"Файл документа недоступен перед обработкой: {source}",
                        }
                    )
                continue
            target = _unique_path(staging_dir / source.name)
            shutil.copy2(source, target)
            staged_files.append(target)
        except Exception as e:
            if issues is not None:
                issues.append(
                    {
                        "severity": "critical",
                        "registry": registry,
                        "file": source.name,
                        "message": f"Не удалось подготовить файл документа к обработке: {e}",
                    }
                )

    if output_dir.exists():
        shutil.rmtree(output_dir, ignore_errors=True)
    output_dir.mkdir(parents=True, exist_ok=True)

    try:
        queue = staged_files
        seen: set[Path] = set()
        archive_count = 0

        while queue and len(seen) < MAX_DOCUMENT_FILES:
            source = queue.pop(0).resolve()
            if source in seen or not source.is_file():
                continue
            seen.add(source)

            is_inside_output = _is_relative_to(source, output_dir)

            if _is_archive(source):
                archive_count += 1
                archive_source = source
                if not is_inside_output:
                    archive_copy = _unique_path(output_dir / source.name)
                    archive_copy.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(source, archive_copy)
                    archive_source = archive_copy.resolve()
                extract_parent = archive_source.parent
                extract_dir = _unique_path(extract_parent / f"{archive_source.stem}_разархивировано")
                if progress:
                    progress(f"Распаковываю архив: {archive_source.name}")
                try:
                    _extract_archive(archive_source, extract_dir)
                    queue.extend(_walk_files([extract_dir]))
                except Exception as e:
                    marker = extract_dir / "ОШИБКА_РАСПАКОВКИ.txt"
                    extract_dir.mkdir(parents=True, exist_ok=True)
                    marker.write_text(
                        f"Исходный архив сохранён здесь: {archive_source}\n"
                        f"Не удалось распаковать {archive_source.name}: {e}",
                        encoding="utf-8",
                    )
                    if issues is not None:
                        issues.append(
                            {
                                "severity": "critical",
                                "registry": registry,
                                "file": archive_source.name,
                                "message": f"Не удалось распаковать архив: {e}",
                            }
                        )
                continue

            if not is_inside_output:
                target = _unique_path(output_dir / source.name)
                target.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(source, target)
    finally:
        shutil.rmtree(staging_dir, ignore_errors=True)

    readable_files = sorted(
        (p for p in output_dir.rglob("*") if p.is_file()),
        key=lambda p: str(p).lower(),
    )
    text = _read_documents_text_from_files(
        readable_files,
        progress=progress,
        issues=issues,
        registry=registry,
    )
    if not readable_files and archive_count == 0:
        text = "[документы не найдены]"
        if issues is not None:
            issues.append(
                {
                    "severity": "important",
                    "registry": registry,
                    "file": "",
                    "message": "Документы не найдены.",
                }
            )
    return text, output_dir


def build_documents_text(files: list[Path], progress=None) -> str:
    """Распаковывает архивы, читает документы и возвращает общий текстовый блок."""
    with tempfile.TemporaryDirectory(prefix="etp_docs_") as tmp:
        text, _ = prepare_documents_for_analysis(files, Path(tmp), progress=progress)
        return text
